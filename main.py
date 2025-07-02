import streamlit as st
import io
import string
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.patches as patches

# --- Main PowerPoint Generation Function ---
def generate_elevation_powerpoint(bay_types_data):
    """
    Creates a professionally styled PowerPoint presentation from bay type data.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    for bay_type in bay_types_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add a title to the slide
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.5))
        p = title_shape.text_frame.paragraphs[0]
        p.text = bay_type['name']
        p.font.name = 'Calibri'
        p.font.bold = True
        p.font.size = Pt(28)

        # --- Corrected Dynamic Scaling and Positioning ---
        # Calculate the true maximum width and height from the user's data
        max_bay_width_cm = max(s['num_bins'] * s['w'] for s in bay_type['shelf_details'].values()) if bay_type['shelf_details'] else 1
        max_bay_height_cm = sum(s['h'] for s in bay_type['shelf_details'].values()) if bay_type['shelf_details'] else 1
        
        # Define the available drawing area on the slide
        drawing_area_width = Inches(8)
        drawing_area_height = Inches(6)

        # Calculate scale to fit the drawing within the area, handling potential division by zero
        scale_w = drawing_area_width / max_bay_width_cm if max_bay_width_cm > 0 else 0
        scale_h = drawing_area_height / max_bay_height_cm if max_bay_height_cm > 0 else 0
        scale = min(scale_w, scale_h) * 0.9 # Use 90% of the calculated scale for better padding

        # Center the drawing horizontally on the slide
        total_drawing_width = max_bay_width_cm * scale
        start_x = Inches(6.67) - (total_drawing_width / 2)
        start_y = Inches(6.8)

        current_y = start_y
        
        # Draw shelves from bottom to top
        for shelf_name in reversed(bay_type['shelves']):
            shelf_info = bay_type['shelf_details'][shelf_name]
            num_bins = shelf_info['num_bins']
            bin_h = shelf_info['h'] * scale
            bin_w = shelf_info['w'] * scale
            shelf_height = bin_h
            shelf_width = num_bins * bin_w

            # 3D effect by drawing a darker, offset shape behind
            shadow_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x + Inches(0.05), current_y - shelf_height + Inches(0.05), shelf_width, shelf_height)
            shadow_shape.fill.solid()
            shadow_shape.fill.fore_color.rgb = RGBColor(180, 180, 180)
            shadow_shape.line.fill.background()

            # Main shelf rectangle
            main_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x, current_y - shelf_height, shelf_width, shelf_height)
            main_shape.fill.solid()
            main_shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            main_shape.line.color.rgb = RGBColor(0, 0, 0)

            # Bin dividers
            for j in range(1, num_bins):
                line_x = start_x + (j * bin_w)
                line = slide.shapes.add_connector(1, line_x, current_y - shelf_height, line_x, current_y) # 1 = straight line
                line.line.color.rgb = RGBColor(120, 120, 120)

            # Shelf Label
            label_box = slide.shapes.add_textbox(start_x - Inches(0.6), current_y - shelf_height, Inches(0.5), shelf_height)
            p = label_box.text_frame.paragraphs[0]
            p.text = shelf_name
            p.font.size = Pt(11)
            p.font.name = 'Calibri'
            
            # Dimension Annotations with connector lines
            dim_x = start_x + shelf_width + Inches(0.5)
            dim_y = current_y - (shelf_height / 2)
            dim_text = f"H: {shelf_info['h']}cm\nW: {shelf_info['w']}cm\nD: {shelf_info['d']}cm"
            dim_box = slide.shapes.add_textbox(dim_x, dim_y - Inches(0.25), Inches(1.5), Inches(0.5))
            tf = dim_box.text_frame
            tf.text = dim_text
            for para in tf.paragraphs:
                para.font.size = Pt(9)
                para.font.name = 'Calibri'
            
            # Connector line from annotation to shelf
            connector = slide.shapes.add_connector(1, dim_x, dim_y, start_x + shelf_width, dim_y)
            connector.line.color.rgb = RGBColor(100, 100, 100)
            connector.line.dash_style = 2 # Dashed line

            current_y -= (shelf_height + Inches(0.08)) # Use a fixed gap for consistency

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# --- Live Preview Function ---
def draw_bay_preview(bay_type_data):
    """
    Generates a static preview of the bay elevation using Matplotlib.
    """
    if not bay_type_data or not bay_type_data['shelves']:
        return None

    fig, ax = plt.subplots(figsize=(10, 6))
    ax.set_aspect('equal')

    current_y = 0
    max_width = 0

    # Draw from bottom up
    for shelf_name in bay_type_data['shelves']:
        shelf_info = bay_type_data['shelf_details'][shelf_name]
        num_bins = shelf_info['num_bins']
        bin_h = shelf_info['h']
        bin_w = shelf_info['w']
        shelf_width = num_bins * bin_w
        max_width = max(max_width, shelf_width)

        # Shelf rectangle
        ax.add_patch(patches.Rectangle((0, current_y), shelf_width, bin_h, fill=True, facecolor='#E0E0E0', edgecolor='black'))
        
        # Bin dividers
        for j in range(1, num_bins):
            ax.plot([j * bin_w, j * bin_w], [current_y, current_y + bin_h], color='gray')
            
        # Shelf label
        ax.text(-5, current_y + bin_h / 2, shelf_name, ha='right', va='center', fontsize=10)

        current_y += bin_h

    ax.set_xlim(-20, max_width + 50)
    ax.set_ylim(0, current_y + 10)
    ax.axis('off')
    
    st.pyplot(fig)

# --- Streamlit App UI ---
st.set_page_config(layout="wide")
st.title("Bay Elevation Generator 📐")
st.markdown("Define your bay types, their shelves, and bin configurations to generate a PowerPoint elevation drawing.")
st.info("ℹ️ **Note:** This application requires `python-pptx` and `matplotlib`. Please ensure they are in your requirements.txt file.")

num_bay_types = st.number_input("How many bay types do you want to define?", min_value=1, max_value=50, value=1, key="num_bay_types")

bay_types_data = []

for i in range(num_bay_types):
    if f"bay_type_name_{i}" not in st.session_state:
        st.session_state[f"bay_type_name_{i}"] = f"Bay Type {i + 1}"

    def update_bay_type_name(idx=i):
        st.session_state[f"bay_type_name_{idx}"] = st.session_state[f"bay_type_name_input_{idx}"]

    header = st.session_state[f"bay_type_name_{i}"].strip() or f"Bay Type {i + 1}"

    with st.expander(header, expanded=True):
        col1, col2 = st.columns([2, 3])
        
        with col1:
            bay_name = st.text_input("Bay Type Name", value=st.session_state[f"bay_type_name_{i}"], key=f"bay_type_name_input_{i}", on_change=update_bay_type_name, args=(i,))
            shelf_count = st.number_input("Number of shelves in this bay?", min_value=1, max_value=26, value=3, key=f"elevation_shelf_count_{i}")
            shelf_names = list(string.ascii_uppercase[:shelf_count])
            
            shelf_details = {}
            
            for shelf_name in shelf_names:
                st.divider()
                st.markdown(f"**Configuration for Shelf {shelf_name}**")
                num_bins = st.number_input("Number of bins in this shelf?", min_value=1, max_value=50, value=5, key=f"num_bins_{i}_{shelf_name}")
                st.markdown(f"**Bin Dimensions for all bins in Shelf {shelf_name} (cm)**")
                c1, c2, c3 = st.columns(3)
                with c1: bin_h = st.number_input("Height", min_value=1.0, value=10.0, key=f"bin_h_{i}_{shelf_name}")
                with c2: bin_w = st.number_input("Width", min_value=1.0, value=10.0, key=f"bin_w_{i}_{shelf_name}")
                with c3: bin_d = st.number_input("Depth", min_value=1.0, value=10.0, key=f"bin_d_{i}_{shelf_name}")
                shelf_details[shelf_name] = {'num_bins': num_bins, 'h': bin_h, 'w': bin_w, 'd': bin_d}

            total_width = 0
            if shelf_details:
                total_width = max(s['num_bins'] * s['w'] for s in shelf_details.values())

            if bay_name.strip():
                bay_types_data.append({
                    "name": bay_name.strip(),
                    "shelves": shelf_names,
                    "shelf_details": shelf_details
                })
        
        with col2:
            st.markdown("**Live Preview**")
            current_bay_data_for_preview = {
                "name": bay_name.strip(),
                "shelves": shelf_names,
                "shelf_details": shelf_details
            }
            draw_bay_preview(current_bay_data_for_preview)
            st.info(f"Calculated Max Bay Width: **{total_width:.2f} cm**")


# --- Generate Button and Download Logic ---
if st.button("Generate PowerPoint Elevation", key="generate_ppt"):
    if any(not bay['name'] for bay in bay_types_data):
        st.error("Please provide a name for all bay types.")
    else:
        with st.spinner("Generating PowerPoint file..."):
            try:
                ppt_buffer = generate_elevation_powerpoint(bay_types_data)
                
                st.success(f"✅ Success! Generated a PowerPoint with {len(bay_types_data)} slides.")
                st.download_button(
                    label="📥 Download PowerPoint File",
                    data=ppt_buffer,
                    file_name="bay_elevations.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"An error occurred during PowerPoint generation: {e}")
